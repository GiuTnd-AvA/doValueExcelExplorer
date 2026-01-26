# =========================
# IMPORT
# =========================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# =========================
# CONFIG
# =========================
OUTPUT_PPTX = r'c:\Users\ciro.andreano\Desktop\doValueExcelExplorer\Migration_Summary_Presentation.pptx'

# =========================
# DATI
# =========================
data = {
    'totale_tabelle_analizzate': 1763,
    'oggetti_estratti': 3070,
    'oggetti_critici': 532,
    'livelli': {
        'L1': 333,
        'L2': 145,
        'L3': 45,
        'L4': 9
    },
    'tipologie': {
        'Stored Procedure': 320,
        'Function': 125,
        'View': 65,
        'Trigger': 22
    },
    'database': {
        'AMS': {'count': 148, 'pct': 44},
        'CORESQL7': {'count': 111, 'pct': 33},
        'S1057': {'count': 37, 'pct': 11},
        'Altri': {'count': 236, 'pct': 12}
    },
    'tabelle_migrare': 760,
    'database_coinvolti': 9,
    'validazione': '100%'
}

# =========================
# FUNZIONI
# =========================

def create_presentation():
    """Crea presentazione PowerPoint con summary migrazione."""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice Blue
    
    # ===== TITOLO =====
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "SQL Server Migration - Executive Summary"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    
    # ===== SUBTITLE =====
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Analisi Criticità e Oggetti da Migrare"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    
    # ===== SEZIONE 1: INPUT =====
    y_pos = 2.5
    
    # Box Input
    input_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1)
    )
    input_box.fill.solid()
    input_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    input_box.line.color.rgb = RGBColor(0, 102, 204)
    input_box.line.width = Pt(2)
    
    input_text = input_box.text_frame
    input_text.text = f"📊 INPUT: {data['totale_tabelle_analizzate']:,} tabelle analizzate → {data['oggetti_estratti']:,} oggetti estratti"
    input_para = input_text.paragraphs[0]
    input_para.font.size = Pt(18)
    input_para.font.bold = True
    input_para.alignment = PP_ALIGN.CENTER
    
    # ===== SEZIONE 2: OGGETTI CRITICI =====
    y_pos = 3.8
    
    # Box Critici
    critici_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(y_pos), Inches(8.4), Inches(1.2)
    )
    critici_box.fill.solid()
    critici_box.fill.fore_color.rgb = RGBColor(255, 245, 230)
    critici_box.line.color.rgb = RGBColor(255, 140, 0)
    critici_box.line.width = Pt(3)
    
    critici_text = critici_box.text_frame
    critici_text.word_wrap = True
    
    p1 = critici_text.paragraphs[0]
    p1.text = f"🎯 OGGETTI CRITICI: {data['oggetti_critici']} (Validati 100%)"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(204, 85, 0)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = critici_text.add_paragraph()
    p2.text = f"L1: {data['livelli']['L1']}  •  L2: {data['livelli']['L2']}  •  L3: {data['livelli']['L3']}  •  L4: {data['livelli']['L4']}"
    p2.font.size = Pt(16)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    
    # ===== SEZIONE 3: TIPOLOGIE OGGETTI =====
    y_pos = 5.1
    
    tipo_header = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(4), Inches(0.4))
    tipo_header_frame = tipo_header.text_frame
    tipo_header_frame.text = "📦 PER TIPOLOGIA"
    tipo_para = tipo_header_frame.paragraphs[0]
    tipo_para.font.size = Pt(14)
    tipo_para.font.bold = True
    tipo_para.alignment = PP_ALIGN.CENTER
    
    # Box tipologie (2x2 grid)
    y_start = 5.5
    x_positions = [0.9, 2.8]
    y_positions = [y_start, y_start + 0.65]
    box_w = 1.7
    box_h = 0.55
    
    tipo_list = [
        ('Stored Proc.', data['tipologie']['Stored Procedure']),
        ('Functions', data['tipologie']['Function']),
        ('Views', data['tipologie']['View']),
        ('Triggers', data['tipologie']['Trigger'])
    ]
    
    tipo_colors = [
        RGBColor(70, 130, 180),   # Steel Blue
        RGBColor(138, 43, 226),   # Blue Violet
        RGBColor(220, 20, 60),    # Crimson
        RGBColor(255, 140, 0)     # Dark Orange
    ]
    
    for idx, (tipo_name, count) in enumerate(tipo_list):
        row = idx // 2
        col = idx % 2
        x_pos = x_positions[col]
        y_pos_tipo = y_positions[row]
        
        tipo_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos), Inches(y_pos_tipo), Inches(box_w), Inches(box_h)
        )
        tipo_box.fill.solid()
        tipo_box.fill.fore_color.rgb = tipo_colors[idx]
        tipo_box.line.color.rgb = RGBColor(255, 255, 255)
        tipo_box.line.width = Pt(1)
        
        tipo_text = tipo_box.text_frame
        tipo_text.word_wrap = True
        tipo_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_tipo = tipo_text.paragraphs[0]
        p_tipo.text = f"{tipo_name}: {count}"
        p_tipo.font.size = Pt(14)
        p_tipo.font.bold = True
        p_tipo.font.color.rgb = RGBColor(255, 255, 255)
        p_tipo.alignment = PP_ALIGN.CENTER
    
    # ===== SEZIONE 4: DISTRIBUZIONE DATABASE =====
    y_pos = 5.1
    
    db_header = slide.shapes.add_textbox(Inches(5.1), Inches(y_pos), Inches(4), Inches(0.4))
    db_header_frame = db_header.text_frame
    db_header_frame.text = "📍 PER DATABASE"
    db_para = db_header_frame.paragraphs[0]
    db_para.font.size = Pt(14)
    db_para.font.bold = True
    db_para.alignment = PP_ALIGN.CENTER
    
    # Database boxes (vertical stack)
    y_pos = 5.5
    x_start = 5.2
    box_width = 1.8
    box_height = 0.35
    spacing = 0.05
    
    db_list = [
        ('AMS', data['database']['AMS']['count'], data['database']['AMS']['pct']),
        ('CORESQL7', data['database']['CORESQL7']['count'], data['database']['CORESQL7']['pct']),
        ('S1057', data['database']['S1057']['count'], data['database']['S1057']['pct']),
        ('Altri 6 DB', data['database']['Altri']['count'], data['database']['Altri']['pct'])
    ]
    
    colors = [
        RGBColor(100, 149, 237),  # Cornflower Blue
        RGBColor(60, 179, 113),   # Medium Sea Green
        RGBColor(255, 165, 0),    # Orange
        RGBColor(169, 169, 169)   # Dark Gray
    ]
    
    for i, (db_name, count, pct) in enumerate(db_list):
        y_pos_db = y_pos + i * (box_height + spacing)
        
        db_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x_start), Inches(y_pos_db), Inches(box_width), Inches(box_height)
        )
        db_box.fill.solid()
        db_box.fill.fore_color.rgb = colors[i]
        db_box.line.color.rgb = RGBColor(0, 0, 0)
        
        db_text = db_box.text_frame
        db_text.word_wrap = True
        db_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p_db = db_text.paragraphs[0]
        p_db.text = f"{db_name}: {count} ({pct}%)"
        p_db.font.size = Pt(12)
        p_db.font.bold = True
        p_db.font.color.rgb = RGBColor(255, 255, 255)
        p_db.alignment = PP_ALIGN.CENTER
    
    # ===== FOOTER =====
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
    footer_frame = footer_box.text_frame
    footer_frame.text = f"✅ Validazione: {data['validazione']} coverage su TOP 1500 oggetti referenced  •  {data['tabelle_migrare']} tabelle da migrare (FASE 0)"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(11)
    footer_para.font.color.rgb = RGBColor(50, 50, 50)
    footer_para.alignment = PP_ALIGN.CENTER
    
    # Salva
    prs.save(OUTPUT_PPTX)
    print(f"\n✅ Presentazione salvata: {OUTPUT_PPTX}")
    print("\nSlide creata con:")
    print(f"  • {data['totale_tabelle_analizzate']:,} tabelle analizzate")
    print(f"  • {data['oggetti_estratti']:,} oggetti estratti")
    print(f"  • {data['oggetti_critici']} oggetti critici validati")
    print(f"  • {data['database_coinvolti']} database coinvolti")
    print(f"  • {data['tabelle_migrare']} tabelle da migrare\n")

# =========================
# MAIN
# =========================

if __name__ == "__main__":
    create_presentation()
